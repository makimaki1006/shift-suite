"""
dash_app.py  – Shift-Suite Dashboard  v1.0 (P0-P3 all-in-one)

◆起動
    python dash_app.py                          # 解析済み out フォルダ 1 件
    python dash_app.py C:\data\out1 C:\out2 ... # 複数施設モード

◆主なタブ
  Overview          : KPI & フォルダ選択
  Heatmap           : ALL ヒートマップ (+ 時間別不足ドリルダウン)
  Role Min-staff    : 職種×method 切替グラフ
  Cost Simulator    : 時給 & 採用人数スライダー → 人件費試算
  Forecast vs Real  : 需要予測線グラフ
  Fairness Tuning   : 夜勤比率スライダー → Jain 指数再計算
  Multi-Facility    : 各施設 KPI 比較表（複数 out フォルダ時のみ）
  PPT Report        : ボタン 1 発で report.pptx を生成（python-pptx）
"""

import dash 
import sys, io, zipfile, json, datetime as dt
from pathlib import Path
from typing import List, Dict

import pandas as pd, numpy as np, plotly.express as px
from dash import Dash, dcc, html, dash_table, Input, Output, State, ctx
import dash_bootstrap_components as dbc

# ppt
from pptx import Presentation
from pptx.util import Inches, Pt

# ────────────────────────────────────────────── 入力フォルダ読み込み
def unpack_if_zip(p: Path) -> Path:
    if p.suffix.lower() != ".zip":
        return p
    tmp = p.parent / (p.stem + "_unzipped")
    if not tmp.exists():
        zipfile.ZipFile(p).extractall(tmp)
    # heat_ALL.xlsx を含むサブフォルダを返す
    for hp in tmp.rglob("heat_ALL.xlsx"):
        return hp.parent
    raise FileNotFoundError("heat_ALL.xlsx not found in ZIP")


def build_dataset(paths: List[Path]) -> Dict[str, Path]:
    """key = facility name (folder名), value = out_dir Path"""
    if not paths:
        p = Path("out")  # デフォルト１件
        return {p.name: p.resolve()}
    return {p.name: unpack_if_zip(p).resolve() for p in paths}


OUT_DIRS = build_dataset([Path(p) for p in sys.argv[1:]])

# ────────────────────────────────────────────── Dash
app = Dash(__name__, external_stylesheets=[dbc.themes.BOOTSTRAP],
           title="Shift-Suite Dashboard", suppress_callback_exceptions=True)

def load_excel(fp: Path, **kw):
    return pd.read_excel(fp, **kw) if fp.exists() else pd.DataFrame()

# ───────────────────────── レイアウト
app.layout = dbc.Container([
    html.H3("Shift-Suite Dashboard"),
    dcc.Dropdown(list(OUT_DIRS.keys()), list(OUT_DIRS.keys())[0],
                 id="facility-dd", style={"maxWidth":300}),
    dcc.Tabs(id="tabs", value="ov", children=[
        dcc.Tab(label="Overview", value="ov"),
        dcc.Tab(label="Heatmap", value="hm"),
        dcc.Tab(label="Role Min-staff", value="min"),
        dcc.Tab(label="Cost Simulator", value="cost"),
        dcc.Tab(label="Forecast vs Real", value="fc"),
        dcc.Tab(label="Fairness Tuning", value="fair"),
        dcc.Tab(label="Multi-Facility", value="multi") if len(OUT_DIRS)>1 else None,
        dcc.Tab(label="PPT Report", value="ppt")
    ]),
    html.Div(id="tab-content")
], fluid=True)

# ───────────────────────── コールバック共通 util
def get_out(fac:str)->Path: return OUT_DIRS[fac]

# ───────────────────────── 各タブ描画
@app.callback(Output("tab-content","children"),
              Input("tabs","value"),
              Input("facility-dd","value"))
def render(tab, fac):
    out = get_out(fac)
    if tab=="ov":   return overview(out)
    if tab=="hm":   return heatmap_tab(out)
    if tab=="min":  return minstaff_tab(out)
    if tab=="cost": return cost_tab(out)
    if tab=="fc":   return forecast_tab(out)
    if tab=="fair": return fairness_tab(out)
    if tab=="multi":return multi_tab()
    if tab=="ppt":  return ppt_tab(out)
    return "N/A"

# ───────────────────────── Overview
def overview(out:Path):
    kpi = load_excel(out/"shortage_role.xlsx")
    lack = int(kpi["lack_h"].sum()) if not kpi.empty else 0

    from shift_suite.tasks.utils import calculate_jain_index
    fair = load_excel(out/"fairness_after.xlsx", sheet_name="after")
    jain = calculate_jain_index(fair["night_ratio"]) if not fair.empty else "–"

    return dbc.Row([
        dbc.Col(dbc.Card([dbc.CardHeader("Total Lack (h)"),
                          dbc.CardBody(html.H2(lack))])),
        dbc.Col(dbc.Card([dbc.CardHeader("Night Jain"),
                          dbc.CardBody(html.H2(jain))])),
        dbc.Col(html.Small(out, style={"fontSize":"0.7rem"}))
    ])

# ───────────────────────── Heatmap + shortage drill
def heatmap_tab(out:Path):
    heat = load_excel(out/"heat_ALL.xlsx", index_col=0)
    if heat.empty: return html.Div("heat_ALL.xlsx missing")
    fig = px.imshow(heat, aspect="auto", color_continuous_scale="Blues",
                    labels=dict(x="Date", y="Time", color="人数"))
    fig.update_layout(clickmode="event+select")
    return html.Div([
        dcc.Graph(id="heat-hm", figure=fig, style={"height":"70vh"}),
        dcc.Graph(id="shortage-graph")    # drill-down placeholder
    ])

@app.callback(Output("shortage-graph","figure"),
              Input("heat-hm","clickData"),
              State("facility-dd","value"))
def show_shortage(clickData, fac):
    if clickData is None: raise dash.exception.PreventUpdate
    date = clickData["points"][0]["x"]
    lack = load_excel(get_out(fac)/"shortage_time.xlsx", index_col=0)
    if lack.empty: return {}
    fig = px.bar(lack[date], title=f"Lack by time – {date}",
                 labels=dict(value="不足人数", index="Time"))
    return fig

# ───────────────────────── Role min-staff
def minstaff_tab(out:Path):
    roles = sorted({p.stem.replace("min_","") for p in out.glob("min_*.csv")})
    return html.Div([
        dcc.Dropdown(roles, roles[0] if roles else None,
                     id="role-dd", style={"maxWidth":300}),
        dcc.RadioItems(["mean-1s","p25","mode"],"p25",
                       id="method-radio", inline=True),
        dcc.Graph(id="min-role-graph")
    ])

@app.callback(Output("min-role-graph","figure"),
              Input("role-dd","value"),Input("method-radio","value"),
              State("facility-dd","value"))
def show_min(role, method, fac):
    if role is None: return {}
    need = pd.read_csv(get_out(fac)/f"min_{role}.csv", index_col=0)
    fig = px.line(need["need"], title=f"{role}  min-staff ({method})",
                  labels=dict(index="Time", value="人数"))
    return fig

# ───────────────────────── Cost simulator
def cost_tab(out:Path):
    roles = sorted({p.stem.replace("min_","") for p in out.glob("min_*.csv")})
    return html.Div([
        html.Div([
            html.Label("対象職種"), dcc.Dropdown(roles, roles[:2], multi=True, id="cost-roles"),
            html.Label("平均時給 (¥)"), dcc.Slider(1000, 3000, 50, value=1500, id="wage"),
            html.Label("採用人数追加"), dcc.Slider(0, 10, 1, value=0, id="hire"),
        ], style={"maxWidth":400}),
        html.Hr(),
        html.H4(id="cost-result")
    ])

@app.callback(Output("cost-result","children"),
              Input("cost-roles","value"),
              Input("wage","value"),
              Input("hire","value"),
              State("facility-dd","value"))
def calc_cost(roles, wage, hire, fac):
    if not roles: return "職種未選択"
    out = get_out(fac)
    lack = load_excel(out/"shortage_role.xlsx")
    lack = lack[lack["role"].isin(roles)] if not lack.empty else pd.DataFrame()
    lack_h = lack["lack_h"].sum() if not lack.empty else 0
    hrs = lack_h - hire*160          # 1 人あたり月 160h 埋める想定
    hrs = max(hrs, 0)
    cost = hrs*wage
    return f"不足 {hrs:.1f}h を派遣で埋めると概算 ¥{cost:,.0f}"

# ───────────────────────── Forecast vs Real
def forecast_tab(out:Path):
    fc = load_excel(out/"forecast.xlsx")
    if fc.empty: return html.Div("forecast.xlsx missing")
    heat = load_excel(out/"heat_ALL.xlsx", index_col=0)
    real = heat.sum()
    df = pd.DataFrame({"ds":pd.to_datetime(fc["ds"]),
                       "Forecast":fc["yhat"],
                       "Real":real.values})
    fig = px.line(df, x="ds", y=["Forecast","Real"],
                  labels=dict(value="人数", ds="Date"))
    return dcc.Graph(figure=fig)

# ───────────────────────── Fairness tuning
def fairness_tab(out:Path):
    base = load_excel(out/"fairness_after.xlsx", sheet_name="before")
    if base.empty: return html.Div("fairness_after.xlsx missing")
    return html.Div([
        dcc.Slider(0,0.5,0.05,value=0.2,id="fair-max"),
        html.Div(id="fair-jain")
    ])

@app.callback(Output("fair-jain","children"),
              Input("fair-max","value"),
              State("facility-dd","value"))
def tune_fair(max_ratio, fac):
    from shift_suite.tasks.utils import calculate_jain_index
    df = load_excel(get_out(fac)/"fairness_after.xlsx", sheet_name="after")
    if df.empty: return "データ無し"
    adj = df["night_ratio"].clip(upper=max_ratio)
    jain = calculate_jain_index(adj)
    return f"許容夜勤比率 {max_ratio:.2f} → Jain = {jain:.3f}"

# ───────────────────────── Multi-facility KPI
def multi_tab():
    records=[]
    for fac,out in OUT_DIRS.items():
        lack = load_excel(out/"shortage_role.xlsx")
        lack_h = lack["lack_h"].sum() if not lack.empty else 0
        from shift_suite.tasks.utils import calculate_jain_index
        fair = load_excel(out/"fairness_after.xlsx", sheet_name="after")
        jain = calculate_jain_index(fair["night_ratio"]) if not fair.empty else None
        records.append({"facility":fac,"lack_h":lack_h,"jain":jain})
    return dash_table.DataTable(records,[{"name":c,"id":c} for c in records[0]],
                                sort_action="native")

# ───────────────────────── PPT report
def ppt_tab(out:Path):
    return html.Div([
        dbc.Button("Generate PPT", id="ppt-btn", n_clicks=0),
        html.Div(id="ppt-msg")
    ])

@app.callback(Output("ppt-msg","children"),
              Input("ppt-btn","n_clicks"),
              State("facility-dd","value"))
def gen_ppt(n, fac):
    if n==0: return ""
    out=get_out(fac)
    prs=Presentation()
    slide=prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text=f"Shift-Suite Report  ({fac})  {dt.date.today()}"
    # KPI
    lack = load_excel(out/"shortage_role.xlsx"); lack_h=int(lack["lack_h"].sum()) if not lack.empty else 0
    slide.shapes.add_textbox(Inches(0.5),Inches(2),Inches(5),Inches(1)).text=f"Lack hours: {lack_h}"

    fp=out/"report.pptx"; prs.save(fp)
    return dbc.Alert(f"PPT saved → {fp}", color="success", duration=4000)

# ───────────────────────── メイン
if __name__ == "__main__":
    app.run(debug=True, port=8502)

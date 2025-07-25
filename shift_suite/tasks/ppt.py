"""Generate a simple PowerPoint summary of analysis results."""

from __future__ import annotations

import logging
import tempfile
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.util import Inches

from shift_suite.i18n import translate as _

log = logging.getLogger(__name__)


def _add_shortage_slide(prs: Presentation, shortage_fp: Path) -> None:
    """Add a bar chart slide summarising shortage by role."""
    import matplotlib.pyplot as plt

    df = pd.read_parquet(shortage_fp)  # read_excel から read_parquet に変更
    if df.empty:
        return

    fig, ax = plt.subplots(figsize=(6, 4))
    df.plot.bar(x="role", y="lack_h", ax=ax)
    ax.set_title(_("Shortage by Role"))
    ax.set_ylabel(_("Shortage Hours"))
    plt.tight_layout()

    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp_png:
        fig.savefig(tmp_png.name)
        plt.close(fig)
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = _("Shortage by Role")
        slide.shapes.add_picture(tmp_png.name, Inches(1), Inches(1), width=Inches(8))
    Path(tmp_png.name).unlink(missing_ok=True)


def _add_cost_slide(prs: Presentation, cost_fp: Path) -> None:
    """Add a slide showing cost simulations."""
    import matplotlib.pyplot as plt

    df = pd.read_parquet(cost_fp)  # read_excel から read_parquet に変更
    if df.empty:
        return

    cost_col = "Cost_Million" if "Cost_Million" in df.columns else df.columns[0]

    fig, ax = plt.subplots(figsize=(6, 4))
    df[cost_col].plot.bar(ax=ax)
    ax.set_ylabel(cost_col)
    ax.set_title(_("Cost Benefit Scenarios"))
    plt.tight_layout()

    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp_png:
        fig.savefig(tmp_png.name)
        plt.close(fig)
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = _("Cost Benefit Scenarios")
        slide.shapes.add_picture(tmp_png.name, Inches(1), Inches(1), width=Inches(8))
    Path(tmp_png.name).unlink(missing_ok=True)


def build_ppt(out_dir: Path) -> Path:
    """Build a PowerPoint report using analysis outputs in *out_dir*."""

    out_dir = Path(out_dir)
    ppt_fp = out_dir / "ShiftSuite_Report.pptx"

    prs = Presentation()

    # title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "ShiftSuite Report"
    if slide.placeholders:
        subtitle = slide.placeholders[1]
        subtitle.text = str(out_dir)

    shortage_fp = (
        out_dir / "shortage_role_summary.parquet"
    )  # ファイル名を .parquet に変更
    cost_fp = out_dir / "cost_benefit.parquet"  # ファイル名を .parquet に変更

    if shortage_fp.exists():
        _add_shortage_slide(prs, shortage_fp)

    if cost_fp.exists():
        _add_cost_slide(prs, cost_fp)

    prs.save(ppt_fp)
    log.info(f"PowerPoint report saved to {ppt_fp}")
    return ppt_fp

# ─────────────────────────────  app.py  (Part 1 / 3)  ──────────────────────────
# Shift-Suite Streamlit GUI + 内蔵ダッシュボード  v1.30.0 (休暇分析機能追加)
# ==============================================================================
# 変更履歴
#   • v1.30.0: 休暇分析機能を追加。leave_analyzer モジュールとの連携を実装。
#   • v1.29.13: st.experimental_rerun() を st.rerun() に修正。
#   • v1.29.12: need_ref_start/end_date_widget の StreamlitAPIException 対策。
#               ファイルアップロード時の日付範囲推定結果を、フラグを用いて
#               次回のスクリプト実行時にウィジェットのデフォルト値として安全に反映するよう修正。
#   • v1.29.11: ログで指摘されたエラー箇所を修正。
#               - shift_sheets_multiselect_widget の StreamlitAPIException 対策。
#               - param_penalty_per_lack の NameError 修正。
#               - progress_bar_exec_main_run 等の NameError 修正。
#               - ログメッセージ内のタイポ修正。
#   • v1.29.10: selectboxのdefault引数エラーをindexに統一し、オプションリストをセッションから正しく参照。
#               全てのウィジェットの値をセッションステートで管理し、初期化を徹底。
#               ヘッダー開始行UIの表示と値の利用を確実化。
#               Excel日付範囲推定の安定化。
#               on_changeコールバックを削除し、よりシンプルなセッションステート管理を目指す。
# ==============================================================================

from __future__ import annotations

import datetime
import io
import json
import logging
import sys
import tempfile
import zipfile
from pathlib import Path
from typing import List, Tuple, Optional, Any

import pandas as pd
import numpy as np
import streamlit as st
from streamlit.runtime import exists as st_runtime_exists
import openpyxl
import plotly.express as px
import plotly.graph_objects as go
import datetime as dt

# ── Shift-Suite task modules ─────────────────────────────────────────────────
from shift_suite.tasks.io_excel import ingest_excel
from shift_suite.tasks.utils import _parse_as_date
from shift_suite.tasks.heatmap import build_heatmap
from shift_suite.tasks.shortage import shortage_and_brief
from shift_suite.tasks.build_stats import build_stats
from shift_suite.tasks.anomaly import detect_anomaly
from shift_suite.tasks.fatigue import train_fatigue
from shift_suite.tasks.cluster import cluster_staff
from shift_suite.tasks.skill_nmf import build_skill_matrix
from shift_suite.tasks.fairness import run_fairness
from shift_suite.tasks.forecast import build_demand_series, forecast_need
from shift_suite.tasks.rl import learn_roster
from shift_suite.tasks.hire_plan import build_hire_plan
from shift_suite.tasks.cost_benefit import analyze_cost_benefit
from shift_suite.tasks.constants import SUMMARY5 as SUMMARY5_CONST
from shift_suite.tasks import leave_analyzer  # ★ 新規インポート
from shift_suite.tasks.leave_analyzer import (
    LEAVE_TYPE_REQUESTED,
    LEAVE_TYPE_PAID,
)
# ──────────────────────────────────────────────────────────────────────────────
from shift_suite.tasks.analyzers import (
    RestTimeAnalyzer,
    WorkPatternAnalyzer,
    AttendanceBehaviorAnalyzer,
    CombinedScoreCalculator,
    LowStaffLoadAnalyzer,
)
from shift_suite.tasks import dashboard

# ── ロガー設定 ─────────────────────────────────
log = logging.getLogger("shift_suite_app")
if not log.handlers:
    log.setLevel(logging.INFO)
    handler = logging.StreamHandler(sys.stdout)
    formatter = logging.Formatter('%(asctime)s [%(levelname)s] %(name)s [%(module)s.%(funcName)s:%(lineno)d] - %(message)s')
    handler.setFormatter(formatter)
    log.addHandler(handler)

    from shift_suite.tasks.utils import log as tasks_log
    if not tasks_log.handlers:
        tasks_log.addHandler(handler)
        tasks_log.setLevel(logging.DEBUG)

# ── 日本語ラベル辞書 & _() ───────────────────────────────────────────────────
JP = {
    "Overview": "概要", "Heatmap": "ヒートマップ", "Shortage": "不足分析",
    "Fatigue": "疲労", "Forecast": "需要予測", "Fairness": "公平性",
    "Cost Sim": "コスト試算", "Hire Plan": "採用計画", "PPT Report": "PPTレポート",
    "Leave Analysis": "休暇分析",  # ★ 追加
    "Slot (min)": "スロット (分)",
    "Need Calculation Settings (Day of Week Pattern)": "📊 Need算出設定 (曜日パターン別)",
    "Reference Period for Need Calculation": "参照期間 (Need算出用)",
    "Rest Time Analysis": "休息時間分析", "Work Pattern Analysis": "勤務パターン分析", "Attendance Analysis": "出勤状況分析", "Combined Score": "総合スコア",
    "Low Staff Load": "少人数勤務分析",
    "Start Date": "開始日", "End Date": "終了日",
    "Statistical Metric for Need": "統計的指標 (Need算出用)",
    "Remove Outliers for Need Calculation": "外れ値を除去してNeedを算出",
    "(Optional) Upper Limit Calculation Method": "(オプション) 上限値算出方法",
    "Min-staff method (for Upper)": "最少人数算出法 (上限値用)",
    "Max-staff method (for Upper)": "最大人数算出法 (上限値用)",
    "Extra modules": "追加モジュール", "Save method": "保存方法",
    "ZIP Download": "ZIP形式でダウンロード", "Save to folder": "フォルダに保存",
    "Run Analysis": "▶ 解析実行", "Cost & Hire Parameters": "💰 コスト・採用計画パラメータ",
    "Standard work hours (h/month)": "所定労働時間 (h/月)",
    "Safety factor (shortage h multiplier)": "安全係数 (不足h上乗せ)",
    "Target coverage rate": "目標充足率",
    "Direct employee labor cost (¥/h)": "正職員 人件費 (¥/h)",
    "Temporary staff labor cost (¥/h)": "派遣 人件費 (¥/h)",
    "One-time hiring cost (¥/person)": "採用一時コスト (¥/人)",
    "Penalty for shortage (¥/h)": "不足ペナルティ (¥/h)",
    "Upload Excel shift file (*.xlsx)": "Excel シフト表 (*.xlsx) をアップロード",
    "Select shift sheets to analyze (multiple)": "解析するシフトシート（複数可）",
    "Header start row (1-indexed)": "ヘッダー開始行 (1-indexed)",
    "File Preview (first 8 rows)": "ファイルプレビュー (先頭8行)",
    "Error during preview display": "プレビューの表示中にエラーが発生しました",
    "Error saving Excel file": "Excelファイルの保存中にエラーが発生しました",
    "Error getting sheet names from Excel": "Excelファイルからシート名の取得中にエラーが発生しました",
    "No analysis target sheets found": "勤務区分シート以外の解析対象シートが見つからないか、全てのシート名に「勤務区分」が含まれています。",
    "Analysis in progress...": "解析準備中...",
    "Ingest: Reading Excel data...": "Excelデータ読み込み中…",
    "Heatmap: Generating heatmap...": "Heatmap生成中…",
    "Shortage: Analyzing shortage...": "Shortage (不足分析) 中…",
    "Stats: Processing...": "Stats (統計情報) 生成中…",
    "Anomaly: Processing...": "Anomaly (異常検知) 中…",
    "Fatigue: Processing...": "Fatigue (疲労分析) 中…",
    "Cluster: Processing...": "Cluster (クラスタリング) 中…",
    "Skill: Processing...": "Skill (スキルNMF) 中…",
    "Fairness: Processing...": "Fairness (公平性分析) 中…",
    "Leave Analysis: Processing...": "Leave Analysis (休暇分析) 中…",  # ★ 追加
    "Need forecast: Processing...": "Need forecast (需要予測) 中…",
    "RL roster (PPO): Processing...": "RL roster (強化学習シフト) 中…",
    "Rest Time Analysis: Processing...": "Rest Time Analysis (休息時間分析) 中…",
    "Work Pattern Analysis: Processing...": "勤務パターン分析 中…",
    "Attendance Analysis: Processing...": "出勤状況分析 中…",
    "Combined Score: Processing...": "総合スコア計算 中…",
    "Low Staff Load: Processing...": "少人数勤務分析 中…",
    "Hire plan: Processing...": "Hire plan (採用計画) 中…",
    "Cost / Benefit: Processing...": "Cost / Benefit (コスト便益分析) 中…",
    "Ingest: Excel data read complete.": "✅ Excelデータ読み込み完了",
    "All processes complete!": "🎉 全ての解析が完了しました！",
    "Error during analysis (ValueError)": "解析中にエラーが発生しました (ValueError)",
    "Required file not found": "必要なファイルが見つかりませんでした",
    "Unexpected error occurred": "予期せぬエラーが発生しました",
    "Save Analysis Results": "📁 解析結果の保存",
    "Output folder": "出力先フォルダ",
    "Open the above path in Explorer.": "エクスプローラーで上記のパスを開いてご確認ください。",
    "Download analysis results as ZIP": "📥 解析結果をZIPでダウンロード",
    "Error creating ZIP file": "ZIPファイルの作成中にエラーが発生しました",
    "Dashboard (Upload ZIP)": "📊 ダッシュボード (ZIP アップロード)",
    "Upload ZIP file of 'out' folder": "out フォルダを ZIP 圧縮してアップロード",
    "heat_ALL.xlsx not found in ZIP": "アップロードされたZIPファイル内に heat_ALL.xlsx が見つかりません。ZIPの構造を確認してください。",
    "Failed to extract ZIP file.": "ZIPファイルの展開に失敗しました。",
    "Uploaded file is not a valid ZIP file.": "アップロードされたファイルは有効なZIPファイルではありません。",
    "Error during ZIP file extraction": "ZIPファイルの展開中にエラーが発生しました",
    "Display Mode": "表示モード", "Raw Count": "人数", "Ratio (staff ÷ need)": "Ratio (staff ÷ need)",
    "Color Scale Max (zmax)": "カラースケール上限 (zmax)",
    "Shortage by Role (hours)": "職種別不足時間 (h)",
    "Shortage by Time (count per day)": "時間帯別不足人数 (日別)",
    "Select date to display": "表示する日付を選択",
    "No date columns in shortage data.": "不足時間データに日付列がありません。",
    "Display all time-slot shortage data": "全時間帯別不足データ表示",
    "Fatigue Score per Staff": "スタッフ別疲労スコア",
    "Demand Forecast (yhat)": "需要予測結果 (yhat)", "Actual (y)": "実績 (y)",
    "Display forecast data": "予測データ表示",
    "Fairness (Night Shift Ratio)": "公平性 (夜勤比率)",
    "Cost Simulation (Million ¥)": "コスト試算 (百万円)",
    "Hiring Plan (Needed FTE)": "採用計画 (必要採用人数)",
    "Hiring Plan Parameters": "採用計画パラメータ",
    "Generate PowerPoint Report (β)": "📊 PowerPointレポート生成 (β版)",
    "Generating PowerPoint report...": "PowerPointレポートを生成中です... 少々お待ちください。",
    "PowerPoint report ready.": "PowerPointレポートの準備ができました。",
    "Download Report (PPTX)": "📥 レポート(PPTX)をダウンロード",
    "python-pptx library required for PPT": "PowerPointレポートの生成には `python-pptx` ライブラリが必要です。\nインストールしてください: `pip install python-pptx`",
    "Error generating PowerPoint report": "PowerPointレポートの生成中にエラーが発生しました",
    "Click button to generate report.": "ボタンをクリックすると、主要な分析結果を含むPowerPointレポートが生成されます（現在はβ版です）。",
}
def _(text_key: str) -> str:
    return JP.get(text_key, text_key)

st.set_page_config(page_title="Shift-Suite", layout="wide", initial_sidebar_state="expanded")
st.title("🗂️ Shift-Suite : 勤務シフト分析ツール")

master_sheet_keyword = "勤務区分"

# --- セッションステートの初期化 (一度だけ実行) ---
if "app_initialized" not in st.session_state:
    st.session_state.app_initialized = True
    st.session_state.analysis_done = False
    st.session_state.work_root_path_str = None
    st.session_state.out_dir_path_str = None
    st.session_state.current_step_for_progress = 0

    today_val = datetime.date.today()

    # サイドバーのウィジェットのキーとデフォルト値をセッションステートに初期設定
    st.session_state.slot_input_widget = 30
    st.session_state.header_row_input_widget = 3
    st.session_state.candidate_sheet_list_for_ui = []
    st.session_state.shift_sheets_multiselect_widget = []
    st.session_state._force_update_multiselect_flag = False

    st.session_state.need_ref_start_date_widget = today_val - datetime.timedelta(days=59) # 初期デフォルト
    st.session_state.need_ref_end_date_widget = today_val - datetime.timedelta(days=1)   # 初期デフォルト
    st.session_state._force_update_need_ref_dates_flag = False
    st.session_state._intended_need_ref_start_date = None
    st.session_state._intended_need_ref_end_date = None

    st.session_state.need_stat_method_options_widget = ["10パーセンタイル", "25パーセンタイル", "中央値", "平均値"]
    st.session_state.need_stat_method_widget = "中央値"
    st.session_state.need_remove_outliers_widget = True

    st.session_state.min_method_for_upper_options_widget = ["mean-1s", "p25", "mode"]
    st.session_state.min_method_for_upper_widget = "p25"
    st.session_state.max_method_for_upper_options_widget = ["mean+1s", "p75"]
    st.session_state.max_method_for_upper_widget = "p75"

    # ★ 休暇分析を含む追加モジュールリスト
    st.session_state.available_ext_opts_widget = [
        "Stats", "Anomaly", "Fatigue", "Cluster", "Skill", "Fairness", "Rest Time Analysis", "Work Pattern Analysis", "Attendance Analysis", "Combined Score", "Low Staff Load", _("Leave Analysis"), "Need forecast", "RL roster (PPO)", "Hire plan", "Cost / Benefit"
    ]
    # デフォルトで休暇分析も選択状態にするかはお好みで
    st.session_state.ext_opts_multiselect_widget = st.session_state.available_ext_opts_widget[:] 

    st.session_state.save_mode_selectbox_options_widget = [_("ZIP Download"), _("Save to folder")]
    st.session_state.save_mode_selectbox_widget = _("ZIP Download")

    st.session_state.std_work_hours_widget = 160
    st.session_state.safety_factor_widget = 1.10
    st.session_state.target_coverage_widget = 0.95
    st.session_state.wage_direct_widget = 1500
    st.session_state.wage_temp_widget = 2200
    st.session_state.hiring_cost_once_widget = 180000
    st.session_state.penalty_per_lack_widget = 4000


    # ★ 休暇分析用パラメータの初期化
    st.session_state.leave_analysis_target_types_widget = [LEAVE_TYPE_REQUESTED, LEAVE_TYPE_PAID] # デフォルトで両方
    st.session_state.leave_concentration_threshold_widget = 3 # 希望休集中度閾値のデフォルト

    # ★ 休暇分析結果格納用
    st.session_state.leave_analysis_results = {}
    
    st.session_state.rest_time_results = None
    st.session_state.rest_time_monthly = None
    st.session_state.work_pattern_results = None
    st.session_state.work_pattern_monthly = None
    st.session_state.attendance_results = None
    st.session_state.combined_score_results = None
    st.session_state.low_staff_load_results = None
    st.session_state.uploaded_files_info = {}
    st.session_state.file_options = {}
    st.session_state.analysis_results = {}
    log.info("セッションステートを初期化しました。")

# --- サイドバーのUI要素 ---
with st.sidebar:
    st.header("🛠️ 解析設定")

    st.number_input(_("Slot (min)"), 5, 120,
                    key="slot_input_widget", help="分析の時間間隔（分）")

    st.subheader("📄 シート選択とヘッダー")

    if st.session_state.get("_force_update_multiselect_flag", False):
        new_options = st.session_state.candidate_sheet_list_for_ui
        current_selection = st.session_state.get("shift_sheets_multiselect_widget", [])
        valid_selection = [s for s in current_selection if s in new_options]
        if not valid_selection and new_options:
             st.session_state.shift_sheets_multiselect_widget = new_options[:]
        elif valid_selection:
             st.session_state.shift_sheets_multiselect_widget = valid_selection
        else:
             st.session_state.shift_sheets_multiselect_widget = []
        st.session_state._force_update_multiselect_flag = False

    st.multiselect(
        _("Select shift sheets to analyze (multiple)"),
        options=st.session_state.candidate_sheet_list_for_ui,
        default=st.session_state.shift_sheets_multiselect_widget,
        key="shift_sheets_multiselect_widget"
    )
    st.number_input(
        _("Header start row (1-indexed)"), 1, 20,
        key="header_row_input_widget"
    )

    st.subheader(_("Need Calculation Settings (Day of Week Pattern)"))

    if st.session_state.get("_force_update_need_ref_dates_flag", False):
        if st.session_state.get("_intended_need_ref_start_date"):
            st.session_state.need_ref_start_date_widget = st.session_state._intended_need_ref_start_date
        if st.session_state.get("_intended_need_ref_end_date"):
            st.session_state.need_ref_end_date_widget = st.session_state._intended_need_ref_end_date
        st.session_state._force_update_need_ref_dates_flag = False
        if "_intended_need_ref_start_date" in st.session_state: del st.session_state["_intended_need_ref_start_date"]
        if "_intended_need_ref_end_date" in st.session_state: del st.session_state["_intended_need_ref_end_date"]

    c1_need_ui, c2_need_ui = st.columns(2)
    with c1_need_ui:
        st.date_input(
            _("Start Date"),
            key="need_ref_start_date_widget", help="Need算出の参照期間の開始日"
        )
    with c2_need_ui:
        st.date_input(
            _("End Date"),
            key="need_ref_end_date_widget", help="Need算出の参照期間の終了日"
        )
    st.caption(_("Reference Period for Need Calculation"))

    current_need_stat_method_idx_val = 0
    try:
        current_need_stat_method_idx_val = st.session_state.need_stat_method_options_widget.index(st.session_state.need_stat_method_widget)
    except (ValueError, AttributeError):
        current_need_stat_method_idx_val = 2
    st.selectbox(
        _("Statistical Metric for Need"), options=st.session_state.need_stat_method_options_widget,
        index=current_need_stat_method_idx_val,
        key="need_stat_method_widget", help="曜日別・時間帯別のNeedを算出する際の統計指標"
    )
    st.checkbox(
        _("Remove Outliers for Need Calculation"),
        key="need_remove_outliers_widget", help="IQR法で外れ値を除去してから統計量を計算します"
    )

    with st.expander(_("(Optional) Upper Limit Calculation Method"), expanded=False):
        current_min_method_upper_idx_val = 0
        try: current_min_method_upper_idx_val = st.session_state.min_method_for_upper_options_widget.index(st.session_state.min_method_for_upper_widget)
        except (ValueError, AttributeError): current_min_method_upper_idx_val = 1
        st.selectbox(
            _("Min-staff method (for Upper)"), options=st.session_state.min_method_for_upper_options_widget,
            index=current_min_method_upper_idx_val, key="min_method_for_upper_widget",
            help="（オプション）ヒートマップの『代表的な上限スタッフ数』の算出方法の一部"
        )
        current_max_method_upper_idx_val = 0
        try: current_max_method_upper_idx_val = st.session_state.max_method_for_upper_options_widget.index(st.session_state.max_method_for_upper_widget)
        except (ValueError, AttributeError): current_max_method_upper_idx_val = 0
        st.selectbox(
            _("Max-staff method (for Upper)"), options=st.session_state.max_method_for_upper_options_widget,
            index=current_max_method_upper_idx_val, key="max_method_for_upper_widget",
            help="（オプション）ヒートマップの『代表的な上限スタッフ数』の算出方法"
        )

    st.divider()
    st.subheader("追加分析モジュール")

    st.multiselect(
        _("Extra modules"), st.session_state.available_ext_opts_widget,
        default=st.session_state.ext_opts_multiselect_widget, # 初期値はセッションステートから
        key="ext_opts_multiselect_widget", help="実行する追加の分析モジュールを選択します。"
    )

    # ★ 休暇分析が選択されている場合のみ、関連パラメータ設定UIを表示
    if _("Leave Analysis") in st.session_state.ext_opts_multiselect_widget:
        with st.expander("📊 " + _("Leave Analysis") + " 設定", expanded=True):
            st.multiselect(
                "分析対象の休暇タイプ",
                options=[LEAVE_TYPE_REQUESTED, LEAVE_TYPE_PAID], # 将来的に 'その他休暇' なども追加可能
                key="leave_analysis_target_types_widget",
                help="分析する休暇の種類を選択します。"
            )
            # 希望休が分析対象に含まれている場合のみ閾値設定を表示
            if LEAVE_TYPE_REQUESTED in st.session_state.leave_analysis_target_types_widget:
                st.number_input(
                    "希望休 集中度判定閾値 (人)", 
                    min_value=1, 
                    step=1,
                    key="leave_concentration_threshold_widget",
                    help="同日にこの人数以上の希望休があった場合に「集中」とみなします。"
                )

    current_save_mode_idx_val = 0
    try: current_save_mode_idx_val = st.session_state.save_mode_selectbox_options_widget.index(st.session_state.save_mode_selectbox_widget)
    except (ValueError, AttributeError): current_save_mode_idx_val = 0
    st.selectbox(
        _("Save method"), options=st.session_state.save_mode_selectbox_options_widget,
        index=current_save_mode_idx_val,
        key="save_mode_selectbox_widget", help="解析結果の保存方法を選択します。"
    )

    with st.expander(_("Cost & Hire Parameters")):
        st.number_input(_("Standard work hours (h/month)"),   100, 300, key="std_work_hours_widget")
        st.slider      (_("Safety factor (shortage h multiplier)"), 1.00, 2.00, key="safety_factor_widget")
        st.slider      (_("Target coverage rate"), 0.50, 1.00, key="target_coverage_widget")
        st.number_input(_("Direct employee labor cost (¥/h)"),   500, 10000, key="wage_direct_widget")
        st.number_input(_("Temporary staff labor cost (¥/h)"),   800, 12000, key="wage_temp_widget")
        st.number_input(_("One-time hiring cost (¥/person)"), 0, 1000000, key="hiring_cost_once_widget")
        st.number_input(_("Penalty for shortage (¥/h)"), 0, 20000, key="penalty_per_lack_widget")

# --- メインコンテンツエリア ---
st.header("1. ファイルアップロードと設定")
uploaded_files = st.file_uploader(
    _("Upload Excel shift file (*.xlsx)"),
    type=["xlsx"],
    key="excel_uploader_main_content_area_key",
    help="勤務実績と勤務区分が記載されたExcelファイルをアップロードしてください。",
    accept_multiple_files=True
)

if uploaded_files:
    for uploaded_file in uploaded_files:
        saved_info = st.session_state.uploaded_files_info.get(uploaded_file.name)
        if saved_info is None or saved_info.get("size") != uploaded_file.size:
            if st.session_state.work_root_path_str is None or not Path(st.session_state.work_root_path_str).exists():
                st.session_state.work_root_path_str = tempfile.mkdtemp(prefix="ShiftSuite_")
                log.info(f"新しい一時ディレクトリを作成しました: {st.session_state.work_root_path_str}")

            st.session_state.analysis_done = False
            work_root_on_upload = Path(st.session_state.work_root_path_str)
            excel_path_for_processing = work_root_on_upload / uploaded_file.name

            try:
                with open(excel_path_for_processing, "wb") as f_wb:
                    f_wb.write(uploaded_file.getbuffer())
                log.info(f"アップロードされたExcelファイルを保存しました: {excel_path_for_processing}")
                st.session_state.uploaded_files_info[uploaded_file.name] = {
                    "path": str(excel_path_for_processing),
                    "size": uploaded_file.size,
                }

                # --- NEW: シート名リストを取得してUIを更新 ---
                try:
                    xls = pd.ExcelFile(excel_path_for_processing)
                    candidate_sheets = [s for s in xls.sheet_names if master_sheet_keyword not in s]
                    if candidate_sheets:
                        st.session_state.candidate_sheet_list_for_ui = candidate_sheets
                        st.session_state.shift_sheets_multiselect_widget = candidate_sheets
                    else:
                        st.session_state.candidate_sheet_list_for_ui = []
                        st.session_state.shift_sheets_multiselect_widget = []
                        st.warning(_("No analysis target sheets found"))
                    st.session_state._force_update_multiselect_flag = True
                except Exception as e_get_sheet:
                    st.error(_("Error getting sheet names from Excel") + f": {e_get_sheet}")
            except Exception as e_save_file_process:
                st.error(_("Error saving Excel file") + f": {e_save_file_process}")

# 「解析実行」ボタン
run_button_disabled_status = not st.session_state.uploaded_files_info or \
                               not st.session_state.get("shift_sheets_multiselect_widget", [])
run_button_clicked = st.button(
    _("Run Analysis"), key="run_analysis_button_final_trigger", use_container_width=True, type="primary",
    disabled=run_button_disabled_status
)

# ─────────────────────────────  app.py  (Part 2 / 3)  ──────────────────────────
if run_button_clicked:
    st.session_state.analysis_done = False
    st.session_state.analysis_results = {}

    for file_name, file_info in st.session_state.uploaded_files_info.items():
        st.session_state.current_step_for_progress = 0

        excel_path_to_use = Path(file_info["path"])
        if st.session_state.work_root_path_str is None or not excel_path_to_use.exists():
            st.error("Excelファイルが正しくアップロードされていません。ファイルを再アップロードしてください。")
            st.stop()

        work_root_exec = excel_path_to_use.parent
        st.session_state.out_dir_path_str = str(work_root_exec / "out")
        out_dir_exec = Path(st.session_state.out_dir_path_str)
        out_dir_exec.mkdir(parents=True, exist_ok=True)
        log.info(f"解析出力ディレクトリ: {out_dir_exec} (file: {file_name})")

        # --- 実行時のUIの値をセッションステートから取得 ---
        param_selected_sheets = st.session_state.shift_sheets_multiselect_widget
        param_header_row = st.session_state.header_row_input_widget
        param_slot = st.session_state.slot_input_widget
        param_ref_start = st.session_state.need_ref_start_date_widget
        param_ref_end = st.session_state.need_ref_end_date_widget
        param_need_stat = st.session_state.need_stat_method_widget
        param_need_outlier = st.session_state.need_remove_outliers_widget
        param_min_method_upper = st.session_state.min_method_for_upper_widget
        param_max_method_upper = st.session_state.max_method_for_upper_widget
        param_ext_opts = st.session_state.ext_opts_multiselect_widget
        param_save_mode = st.session_state.save_mode_selectbox_widget
        param_std_work_hours = st.session_state.std_work_hours_widget
        param_safety_factor = st.session_state.safety_factor_widget
        param_target_coverage = st.session_state.target_coverage_widget
        param_wage_direct = st.session_state.wage_direct_widget
        param_wage_temp = st.session_state.wage_temp_widget
        param_hiring_cost = st.session_state.hiring_cost_once_widget
        param_penalty_lack = st.session_state.penalty_per_lack_widget
        
        # ★ 休暇分析用パラメータの取得
        param_leave_target_types = st.session_state.leave_analysis_target_types_widget
        param_leave_concentration_threshold = st.session_state.leave_concentration_threshold_widget
        
        # ★ セッションステート内の前回結果をクリア
        st.session_state.leave_analysis_results = {}
        # --- UI値取得ここまで ---
    
        st.session_state.rest_time_results = None
        st.session_state.work_pattern_results = None
        st.session_state.attendance_results = None
        st.session_state.combined_score_results = None
        st.session_state.low_staff_load_results = None
        progress_text_area = st.empty()
        progress_bar_val = st.progress(0)
        total_steps_exec_run = 3 + len(param_ext_opts)
    
        def update_progress_exec_run(step_name_key_exec: str):
            st.session_state.current_step_for_progress += 1
            progress_percentage_exec_run = int((st.session_state.current_step_for_progress / total_steps_exec_run) * 100)
            progress_percentage_exec_run = min(progress_percentage_exec_run, 100)
            try:
                progress_bar_val.progress(progress_percentage_exec_run)
                progress_text_area.info(f"⚙️ {st.session_state.current_step_for_progress}/{total_steps_exec_run} - {_(step_name_key_exec)}")
            except Exception as e_prog_exec_run: 
                log.warning(f"進捗表示の更新中にエラー: {e_prog_exec_run}")
    
        st.markdown("---")
        st.header("2. 解析処理")
        try:
            if param_selected_sheets and excel_path_to_use:
                update_progress_exec_run("File Preview (first 8 rows)")
                st.subheader(_("File Preview (first 8 rows)"))
                try:
                    preview_df_exec_run = pd.read_excel(excel_path_to_use, sheet_name=param_selected_sheets[0], header=None, nrows=8)
                    st.dataframe(preview_df_exec_run.astype(str), use_container_width=True)
                except Exception as e_prev_exec_run: 
                    st.warning(_("Error during preview display") + f": {e_prev_exec_run}")
                    log.warning(f"プレビュー表示エラー: {e_prev_exec_run}", exc_info=True)
            else: 
                st.warning("プレビューを表示するシートが選択されていないか、ファイルパスが無効です。")
    
            update_progress_exec_run("Ingest: Reading Excel data...")
            long_df, wt_df = ingest_excel(excel_path_to_use, shift_sheets=param_selected_sheets, header_row=param_header_row)
            log.info(f"Ingest完了. long_df shape: {long_df.shape}, wt_df shape: {wt_df.shape if wt_df is not None else 'N/A'}")
            st.success(_("Ingest: Excel data read complete."))
    
            update_progress_exec_run("Heatmap: Generating heatmap...")
            build_heatmap(
                long_df, out_dir_exec, param_slot,
                ref_start_date_for_need=param_ref_start,
                ref_end_date_for_need=param_ref_end,
                need_statistic_method=param_need_stat,
                need_remove_outliers=param_need_outlier,
                need_iqr_multiplier=1.5,
                min_method=param_min_method_upper,
                max_method=param_max_method_upper
            )
            st.success("✅ Heatmap生成完了")
    
            update_progress_exec_run("Shortage: Analyzing shortage...")
            shortage_result_exec_run = shortage_and_brief(out_dir_exec, param_slot)
            if shortage_result_exec_run is None: 
                st.warning("Shortage (不足分析) の一部または全てが完了しませんでした。")
            else: 
                st.success("✅ Shortage (不足分析) 完了")
    
            # ★----- 休暇分析モジュールの実行 -----★
            # "休暇分析" (日本語) が選択されているか確認
            if _("Leave Analysis") in param_ext_opts:
                update_progress_exec_run("Leave Analysis: Processing...")
                st.info(f"{_('Leave Analysis')} 処理中…")
                try:
                    if 'long_df' in locals() and not long_df.empty:
                        # 1. 日次・職員別の休暇取得フラグデータを生成
                        daily_leave_df = leave_analyzer.get_daily_leave_counts(
                            long_df,
                            target_leave_types=param_leave_target_types
                        )
                        st.session_state.leave_analysis_results['daily_leave_df'] = daily_leave_df
                        
                        if not daily_leave_df.empty:
                            leave_results_temp = {} # 一時的な結果格納用
                            
                            # 2. 希望休関連の集計と分析
                            if LEAVE_TYPE_REQUESTED in param_leave_target_types:
                                requested_leave_daily = daily_leave_df[daily_leave_df['leave_type'] == LEAVE_TYPE_REQUESTED]
                                if not requested_leave_daily.empty:
                                    leave_results_temp['summary_dow_requested'] = leave_analyzer.summarize_leave_by_day_count(requested_leave_daily.copy(), period='dayofweek')
                                    leave_results_temp['summary_month_period_requested'] = leave_analyzer.summarize_leave_by_day_count(requested_leave_daily.copy(), period='month_period')
                                    leave_results_temp['summary_month_requested'] = leave_analyzer.summarize_leave_by_day_count(requested_leave_daily.copy(), period='month')
                                    
                                    daily_requested_applicants_counts = leave_analyzer.summarize_leave_by_day_count(requested_leave_daily.copy(), period='date')
                                    leave_results_temp['concentration_requested'] = leave_analyzer.analyze_leave_concentration(
                                        daily_requested_applicants_counts,
                                        leave_type_to_analyze=LEAVE_TYPE_REQUESTED,
                                        concentration_threshold=param_leave_concentration_threshold
                                    )
                                    # --- 新規: 勤務予定人数との比較データ作成 ---
                                    try:
                                        total_staff_per_day = (
                                            long_df[long_df["parsed_slots_count"] > 0]
                                            .assign(date=lambda df: pd.to_datetime(df["ds"]).dt.normalize())
                                            .groupby("date")["staff"].nunique()
                                            .reset_index(name="total_staff")
                                        )
                                        staff_balance = total_staff_per_day.merge(
                                            daily_requested_applicants_counts.rename(columns={"total_leave_days": "leave_applicants_count"})[["date", "leave_applicants_count"]],
                                            on="date",
                                            how="left",
                                        )
                                        staff_balance["leave_applicants_count"] = staff_balance["leave_applicants_count"].fillna(0).astype(int)
                                        staff_balance["non_leave_staff"] = staff_balance["total_staff"] - staff_balance["leave_applicants_count"]
                                        leave_results_temp["staff_balance_daily"] = staff_balance
                                    except Exception as e:
                                        log.error(f"勤務予定人数の計算中にエラー: {e}")
                                else:
                                    log.info(f"{LEAVE_TYPE_REQUESTED} のデータが見つからなかったため、関連する集計・分析をスキップしました。")
                                    leave_results_temp['summary_dow_requested'] = pd.DataFrame()
                                    leave_results_temp['summary_month_period_requested'] = pd.DataFrame()
                                    leave_results_temp['summary_month_requested'] = pd.DataFrame()
                                    leave_results_temp['concentration_requested'] = pd.DataFrame()
                            
                            # 3. 有給休暇関連の集計
                            if LEAVE_TYPE_PAID in param_leave_target_types:
                                paid_leave_daily = daily_leave_df[daily_leave_df['leave_type'] == LEAVE_TYPE_PAID]
                                if not paid_leave_daily.empty:
                                    leave_results_temp['summary_dow_paid'] = leave_analyzer.summarize_leave_by_day_count(paid_leave_daily.copy(), period='dayofweek')
                                    leave_results_temp['summary_month_paid'] = leave_analyzer.summarize_leave_by_day_count(paid_leave_daily.copy(), period='month')
                                else:
                                    log.info(f"{LEAVE_TYPE_PAID} のデータが見つからなかったため、関連する集計をスキップしました。")
                                    leave_results_temp['summary_dow_paid'] = pd.DataFrame()
                                    leave_results_temp['summary_month_paid'] = pd.DataFrame()
                            
                            # 4. 職員別休暇リスト (終日のみ)
                            leave_results_temp['staff_leave_list'] = leave_analyzer.get_staff_leave_list(long_df, target_leave_types=param_leave_target_types)
                            
                            st.session_state.leave_analysis_results.update(leave_results_temp)
                            st.success(f"✅ {_('Leave Analysis')} 完了")
                        else:
                            st.info(f"{_('Leave Analysis')}: 分析対象となる休暇データが見つかりませんでした。")
                    else:
                        st.warning(f"{_('Leave Analysis')}: 前提となる long_df が存在しないか空のため、処理をスキップしました。")
                except Exception as e_leave:
                    st.error(f"{_('Leave Analysis')} の処理中にエラーが発生しました: {e_leave}")
                    log.error(f"休暇分析エラー: {e_leave}", exc_info=True)
            # ★----- 休暇分析モジュールの実行ここまで -----★
    
            # 他の追加モジュールの実行
            for opt_module_name_exec_run in st.session_state.available_ext_opts_widget:
                if opt_module_name_exec_run in param_ext_opts and opt_module_name_exec_run != _("Leave Analysis"):
                    progress_key_exec_run = f"{opt_module_name_exec_run}: Processing..."
                    update_progress_exec_run(progress_key_exec_run)
                    st.info(f"{_(opt_module_name_exec_run)} 処理中…")
                    try:
                        if opt_module_name_exec_run == "Stats": 
                            build_stats(out_dir_exec)
                        elif opt_module_name_exec_run == "Anomaly": 
                            detect_anomaly(out_dir_exec)
                        elif opt_module_name_exec_run == "Fatigue": 
                            train_fatigue(long_df, out_dir_exec)
                        elif opt_module_name_exec_run == "Cluster": 
                            cluster_staff(long_df, out_dir_exec)
                        elif opt_module_name_exec_run == "Skill": 
                            build_skill_matrix(long_df, out_dir_exec)
                        elif opt_module_name_exec_run == "Fairness": 
                            run_fairness(long_df, out_dir_exec)
                        elif opt_module_name_exec_run == "Rest Time Analysis":
                            rta = RestTimeAnalyzer()
                            st.session_state.rest_time_results = rta.analyze(long_df)
                            st.session_state.rest_time_results.to_csv(out_dir_exec / "rest_time.csv", index=False)
                            st.session_state.rest_time_monthly = rta.monthly(st.session_state.rest_time_results)
                            if st.session_state.rest_time_monthly is not None:
                                st.session_state.rest_time_monthly.to_csv(out_dir_exec / "rest_time_monthly.csv", index=False)
                        elif opt_module_name_exec_run == "Work Pattern Analysis":
                            wpa = WorkPatternAnalyzer()
                            st.session_state.work_pattern_results = wpa.analyze(long_df)
                            st.session_state.work_pattern_results.to_csv(out_dir_exec / "work_patterns.csv", index=False)
                            st.session_state.work_pattern_monthly = wpa.analyze_monthly(long_df)
                            if st.session_state.work_pattern_monthly is not None:
                                st.session_state.work_pattern_monthly.to_csv(out_dir_exec / "work_pattern_monthly.csv", index=False)
                        elif opt_module_name_exec_run == "Attendance Analysis":
                            st.session_state.attendance_results = AttendanceBehaviorAnalyzer().analyze(long_df)
                            st.session_state.attendance_results.to_csv(out_dir_exec / "attendance.csv", index=False)
                        elif opt_module_name_exec_run == "Low Staff Load":
                            lsl = LowStaffLoadAnalyzer()
                            st.session_state.low_staff_load_results = lsl.analyze(long_df, threshold=0.25)
                            st.session_state.low_staff_load_results.to_csv(out_dir_exec / "low_staff_load.csv", index=False)
                        elif opt_module_name_exec_run == "Combined Score":
                            rest_df = st.session_state.rest_time_results if st.session_state.rest_time_results is not None else pd.DataFrame()
                            work_df = st.session_state.work_pattern_results if st.session_state.work_pattern_results is not None else pd.DataFrame()
                            att_df = st.session_state.attendance_results if st.session_state.attendance_results is not None else pd.DataFrame()
                            st.session_state.combined_score_results = CombinedScoreCalculator().calculate(rest_df, work_df, att_df)
                            st.session_state.combined_score_results.to_csv(out_dir_exec / "combined_score.csv", index=False)
                        elif opt_module_name_exec_run == "Need forecast":
                            demand_csv_exec_run_fc = out_dir_exec / "demand_series.csv"
                            forecast_xls_exec_run_fc = out_dir_exec / "forecast.xlsx"
                            heat_all_for_fc_exec_run_fc = out_dir_exec / "heat_ALL.xlsx"
                            if not heat_all_for_fc_exec_run_fc.exists(): 
                                st.warning(f"Need forecast: 必須ファイル {heat_all_for_fc_exec_run_fc.name} が見つかりません。")
                            else:
                                build_demand_series(heat_all_for_fc_exec_run_fc, demand_csv_exec_run_fc)
                                if demand_csv_exec_run_fc.exists(): 
                                    forecast_need(demand_csv_exec_run_fc, forecast_xls_exec_run_fc)
                                else: 
                                    st.warning("Need forecast: demand_series.csv の生成に失敗しました。")
                        elif opt_module_name_exec_run == "RL roster (PPO)":
                            demand_csv_rl_exec_run_rl = out_dir_exec / "demand_series.csv"
                            rl_roster_xls_exec_run_rl = out_dir_exec / "rl_roster.xlsx"
                            if demand_csv_rl_exec_run_rl.exists(): 
                                learn_roster(demand_csv_rl_exec_run_rl, rl_roster_xls_exec_run_rl)
                            else: 
                                st.warning("RL Roster: 需要予測データ (demand_series.csv) がありません。")
                        elif opt_module_name_exec_run == "Hire plan":
                            demand_csv_hp_exec_run_hp = out_dir_exec / "demand_series.csv"
                            hire_xls_exec_run_hp = out_dir_exec / "hire_plan.xlsx"
                            if demand_csv_hp_exec_run_hp.exists(): 
                                build_hire_plan(demand_csv_hp_exec_run_hp, hire_xls_exec_run_hp, param_std_work_hours, param_safety_factor, param_target_coverage)
                            else: 
                                st.warning("Hire Plan: 需要予測データ (demand_series.csv) がありません。")
                        elif opt_module_name_exec_run == "Cost / Benefit":
                            analyze_cost_benefit(out_dir_exec, param_wage_direct, param_wage_temp, param_hiring_cost, param_penalty_lack)
                        st.success(f"✅ {_(opt_module_name_exec_run)} 完了")
                    except FileNotFoundError as fe_opt_exec_run_loop:
                        st.error(f"{_(opt_module_name_exec_run)} の処理中にエラー (ファイル未検出): {fe_opt_exec_run_loop}")
                        log.error(f"{opt_module_name_exec_run} 処理エラー (FileNotFoundError): {fe_opt_exec_run_loop}", exc_info=True)
                    except Exception as e_opt_exec_run_loop:
                        st.error(f"{_(opt_module_name_exec_run)} の処理中にエラーが発生しました: {e_opt_exec_run_loop}")
                        log.error(f"{opt_module_name_exec_run} 処理エラー: {e_opt_exec_run_loop}", exc_info=True)
    
            progress_bar_val.progress(100)
            progress_text_area.success("✨ 全工程完了！")
            st.balloons()
            st.success(_("All processes complete!"))
            st.session_state.analysis_done = True
        except ValueError as ve_exec_run_main:
            st.error(_("Error during analysis (ValueError)") + f": {ve_exec_run_main}")
            log.error(f"解析エラー (ValueError): {ve_exec_run_main}", exc_info=True)
            st.session_state.analysis_done = False
        except FileNotFoundError as fe_exec_run_main:
            st.error(_("Required file not found") + f": {fe_exec_run_main}")
            log.error(f"ファイル未検出エラー: {fe_exec_run_main}", exc_info=True)
            st.session_state.analysis_done = False
        except Exception as e_exec_run_main:
            st.error(_("Unexpected error occurred") + f": {e_exec_run_main}")
            log.error(f"予期せぬエラー: {e_exec_run_main}", exc_info=True)
            st.session_state.analysis_done = False
        finally:
            if 'progress_bar_val' in locals() and progress_bar_val is not None: 
                progress_bar_val.empty()
            if 'progress_text_area' in locals() and progress_text_area is not None: 
                progress_text_area.empty()
    
        if st.session_state.analysis_done and st.session_state.out_dir_path_str:
            out_dir_to_save_exec_main_run = Path(st.session_state.out_dir_path_str)
            if out_dir_to_save_exec_main_run.exists():
                st.markdown("---")
                st.header("3. " + _("Save Analysis Results"))
                current_save_mode_exec_main_run = st.session_state.save_mode_selectbox_widget
                if current_save_mode_exec_main_run == _("Save to folder"):
                    st.info(_("Output folder") + f": `{out_dir_to_save_exec_main_run}`")
                    st.markdown(_("Open the above path in Explorer."))
                else: # ZIP Download
                    zip_base_name_exec_main_run = f"ShiftSuite_Analysis_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}"
                    if st.session_state.work_root_path_str is None:
                        st.error("一時作業ディレクトリが見つかりません。ZIPファイルの作成に失敗しました。")
                    else:
                        work_root_for_zip_dl_exec_main_run = Path(st.session_state.work_root_path_str)
                        zip_path_obj_to_download_exec_main_run = work_root_for_zip_dl_exec_main_run / f"{zip_base_name_exec_main_run}.zip"
                        try:
                            with zipfile.ZipFile(zip_path_obj_to_download_exec_main_run, 'w', zipfile.ZIP_DEFLATED) as zf_dl_exec_main_run:
                                for file_to_zip_dl_exec_main_run in out_dir_to_save_exec_main_run.rglob('*'):
                                    if file_to_zip_dl_exec_main_run.is_file():
                                        zf_dl_exec_main_run.write(file_to_zip_dl_exec_main_run, file_to_zip_dl_exec_main_run.relative_to(out_dir_to_save_exec_main_run))
                            with open(zip_path_obj_to_download_exec_main_run, "rb") as fp_zip_data_to_download_dl_exec_main_run:
                                st.download_button(label=_("Download analysis results as ZIP"), data=fp_zip_data_to_download_dl_exec_main_run,
                                                   file_name=f"{zip_base_name_exec_main_run}.zip", mime="application/zip", use_container_width=True)
                            log.info(f"ZIPファイルを作成し、ダウンロードボタンを表示しました: {zip_path_obj_to_download_exec_main_run}")
                        except Exception as e_zip_final_exec_run_main_ex_v3:
                            st.error(_("Error creating ZIP file") + f": {e_zip_final_exec_run_main_ex_v3}")
                            log.error(f"ZIP作成エラー (最終段階): {e_zip_final_exec_run_main_ex_v3}", exc_info=True)
        else:
            log.warning(f"解析は完了しましたが、出力ディレクトリ '{out_dir_to_save_exec_main_run}' が見つかりません。")

        st.session_state.analysis_results[file_name] = {
            "out_dir_path_str": st.session_state.out_dir_path_str,
            "leave_analysis_results": st.session_state.leave_analysis_results,
            "rest_time_results": st.session_state.rest_time_results,
            "work_pattern_results": st.session_state.work_pattern_results,
            "attendance_results": st.session_state.attendance_results,
            "combined_score_results": st.session_state.combined_score_results,
            "low_staff_load_results": st.session_state.low_staff_load_results,
        }

st.session_state.analysis_done = True

# 完全修正版 - 休暇分析結果表示コード全体

# Plotlyの全体問題を修正した休暇分析コード

# ★ 新しい「休暇分析」タブの表示 (解析が完了し、休暇分析が選択されている場合)
# ─────────────────────────────  app.py  (Part 3 / 3)  ──────────────────────────
def display_overview_tab(tab_container, data_dir):
    with tab_container:
        st.subheader(_("Overview"))
        kpi_fp = data_dir / "shortage_role.xlsx"; lack_h = 0.0
        if kpi_fp.exists():
            try: df_sh_role = pd.read_excel(kpi_fp); lack_h = df_sh_role["lack_h"].sum() if "lack_h" in df_sh_role else 0.0
            except Exception as e: st.warning(f"shortage_role.xlsx 読込/集計エラー: {e}")
        fair_fp_meta = data_dir / "fairness_before.xlsx"; jain_display = "N/A"
        if fair_fp_meta.exists():
            try:
                meta_df = pd.read_excel(fair_fp_meta, sheet_name="meta_summary")
                jain_row = meta_df[meta_df["metric"] == "jain_index"]
                if not jain_row.empty: jain_display = f"{float(jain_row['value'].iloc[0]):.3f}"
            except Exception: pass 
        c1, c2 = st.columns(2)
        c1.metric(_("不足時間(h)"), f"{lack_h:.1f}")
        c2.metric("夜勤 Jain指数", jain_display)

def display_heatmap_tab(tab_container, data_dir):
    with tab_container:
        st.subheader(_("Heatmap"))
        fp = data_dir / "heat_ALL.xlsx"
        if fp.exists():
            try:
                df_heat = pd.read_excel(fp, index_col=0)
                mode_opts = {"Raw": _("Raw Count"), "Ratio": _("Ratio (staff ÷ need)")}
                mode_lbl = st.radio(_("Display Mode"), list(mode_opts.values()), horizontal=True, key="dash_heat_mode_radio")
                mode_key = [k for k,v in mode_opts.items() if v == mode_lbl][0]
                z_def, z_min, z_max, z_stp = (11.,1.,50.,1.) if mode_key=="Raw" else (1.5,.1,3.,.1)
                disp_df_heat = df_heat.drop(columns=[c for c in SUMMARY5_CONST if c in df_heat.columns], errors="ignore")

                if mode_key == "Raw":
                    pos_vals = disp_df_heat[disp_df_heat > 0].stack()
                    p90 = float(pos_vals.quantile(0.90)) if not pos_vals.empty else z_def
                    p95 = float(pos_vals.quantile(0.95)) if not pos_vals.empty else z_def
                    p99 = float(pos_vals.quantile(0.99)) if not pos_vals.empty else z_def
                    zmode = st.selectbox("zmax mode", ["Manual", "90th percentile", "95th percentile", "99th percentile"], key="dash_heat_zmax_mode")
                    if zmode == "Manual":
                        zmax = st.slider(_("Color Scale Max (zmax)"), z_min, z_max, z_def, z_stp, key="dash_heat_zmax_slider")
                    else:
                        perc_map = {"90th percentile": p90, "95th percentile": p95, "99th percentile": p99}
                        z_val = perc_map.get(zmode, z_def)
                        st.slider(_("Color Scale Max (zmax)"), z_min, z_max, z_val, z_stp, key="dash_heat_zmax_slider", disabled=True)
                        zmax = z_val
                    fig = px.imshow(disp_df_heat, aspect="auto", color_continuous_scale="Blues", zmax=zmax, labels={"x":_("Date"),"y":_("Time"),"color":_("Raw Count")})
                else:
                    zmax = st.slider(_("Color Scale Max (zmax)"), z_min, z_max, z_def, z_stp, key="dash_heat_zmax_slider")
                    if "need" in df_heat.columns and "staff" in df_heat.columns and not disp_df_heat.empty :
                        ratio_calc_df = disp_df_heat.copy()
                        # Ratio計算: 各日付列の各時間帯の値を、その時間帯の staff 値 / need 値で更新
                        if 'need' in df_heat.columns and df_heat['need'].replace(0, pd.NA).notna().any(): # needが0でない有効な値を持つか
                            ratio_display_df = disp_df_heat.apply(lambda date_col: date_col / df_heat['need'].replace(0, pd.NA), axis=0)
                            ratio_display_df = ratio_display_df.clip(upper=zmax)
                            fig = px.imshow(ratio_display_df, aspect="auto", color_continuous_scale=px.colors.sequential.RdBu_r, zmin=0, zmax=zmax, labels={"x":_("Date"),"y":_("Time"),"color":_("Ratio (staff ÷ need)")})
                        else:
                            st.warning("Ratio表示に必要な'need'列データが0または存在しません。")
                            fig = go.Figure()
                    else:
                        st.warning("Ratio表示に必要な'staff'列、'need'列、または日付データが見つかりません。")
                        fig = go.Figure()
                st.plotly_chart(fig, use_container_width=True)
            except Exception as e: st.error(f"ヒートマップ表示エラー: {e}")
        else: st.info(_("Heatmap") + " (heat_ALL.xlsx) " + _("が見つかりません。"))

def display_shortage_tab(tab_container, data_dir):
    with tab_container:
        st.subheader(_("Shortage"))
        fp_s_role = data_dir / "shortage_role.xlsx"
        if fp_s_role.exists():
            try:
                df_s_role = pd.read_excel(fp_s_role); st.dataframe(df_s_role,use_container_width=True,hide_index=True)
                if "role" in df_s_role and "lack_h" in df_s_role: st.bar_chart(df_s_role.set_index("role")["lack_h"], color="#FFA500")
            except Exception as e: st.error(f"shortage_role.xlsx 表示エラー: {e}")
        else: st.info(_("Shortage") + " (shortage_role.xlsx) " + _("が見つかりません。"))
        st.markdown("---")
        fp_s_time = data_dir / "shortage_time.xlsx"
        if fp_s_time.exists():
            try:
                df_s_time = pd.read_excel(fp_s_time, index_col=0)
                st.write(_("Shortage by Time (count per day)"))
                avail_dates = df_s_time.columns.tolist()
                if avail_dates:
                    sel_date = st.selectbox(_("Select date to display"), avail_dates, key="dash_short_time_date")
                    if sel_date: st.bar_chart(df_s_time[sel_date])
                else: st.info(_("No date columns in shortage data."))
                with st.expander(_("Display all time-slot shortage data")): st.dataframe(df_s_time, use_container_width=True)
            except Exception as e: st.error(f"shortage_time.xlsx 表示エラー: {e}")
        else: st.info(_("Shortage") + " (shortage_time.xlsx) " + _("が見つかりません。"))
        
def display_fatigue_tab(tab_container, data_dir):
    with tab_container:
        st.subheader(_("Fatigue Score per Staff"))
        fp = data_dir / "fatigue_score.xlsx"
        if fp.exists():
            try:
                df = pd.read_excel(fp) 
                st.dataframe(df, use_container_width=True, hide_index=True)
                if "fatigue_score" in df and "staff" in df: 
                    st.bar_chart(df.set_index("staff")["fatigue_score"])
            except Exception as e: st.error(f"fatigue_score.xlsx 表示エラー: {e}")
        else: st.info(_("Fatigue") + " (fatigue_score.xlsx) " + _("が見つかりません。"))

def display_forecast_tab(tab_container, data_dir):
    with tab_container:
        st.subheader(_("Demand Forecast (yhat)"))
        fp_fc = data_dir / "forecast.xlsx"
        if fp_fc.exists():
            try:
                df_fc = pd.read_excel(fp_fc, parse_dates=["ds"])
                fig = go.Figure()
                if "ds" in df_fc and "yhat" in df_fc:
                    fig.add_trace(go.Scatter(x=df_fc["ds"], y=df_fc["yhat"], mode='lines+markers', name=_("Demand Forecast (yhat)")))
                fp_demand = data_dir / "demand_series.csv"
                if fp_demand.exists():
                    df_actual = pd.read_csv(fp_demand, parse_dates=["ds"])
                    if "ds" in df_actual and "y" in df_actual:
                         fig.add_trace(go.Scatter(x=df_actual["ds"], y=df_actual["y"], mode='lines', name=_("Actual (y)"), line=dict(dash='dash')))
                fig.update_layout(title=_("Demand Forecast vs Actual"), xaxis_title=_("Date"), yaxis_title=_("Demand"))
                st.plotly_chart(fig, use_container_width=True)
                with st.expander(_("Display forecast data")): st.dataframe(df_fc, use_container_width=True, hide_index=True)
            except Exception as e: st.error(f"forecast.xlsx 表示エラー: {e}")
        else: st.info(_("Forecast") + " (forecast.xlsx) " + _("が見つかりません。"))

def display_fairness_tab(tab_container, data_dir):
    with tab_container:
        st.subheader(_("Fairness (Night Shift Ratio)"))
        fp = data_dir / "fairness_after.xlsx"
        if fp.exists():
            try:
                df = pd.read_excel(fp); st.dataframe(df, use_container_width=True, hide_index=True)
                if "staff" in df and "night_ratio" in df: st.bar_chart(df.set_index("staff")["night_ratio"], color="#FF8C00")
            except Exception as e: st.error(f"fairness_after.xlsx 表示エラー: {e}")
        else: st.info(_("Fairness") + " (fairness_after.xlsx) " + _("が見つかりません。"))

def display_costsim_tab(tab_container, data_dir):
    with tab_container:
        st.subheader(_("Cost Simulation (Million ¥)"))
        fp = data_dir / "cost_benefit.xlsx"
        if fp.exists():
            try:
                df = pd.read_excel(fp, index_col=0)
                if "Cost_Million" in df: st.bar_chart(df["Cost_Million"])
                st.dataframe(df, use_container_width=True)
            except Exception as e: st.error(f"cost_benefit.xlsx 表示エラー: {e}")
        else: st.info(_("Cost Sim") + " (cost_benefit.xlsx) " + _("が見つかりません。"))

def display_hireplan_tab(tab_container, data_dir):
    with tab_container:
        st.subheader(_("Hiring Plan (Needed FTE)"))
        fp = data_dir / "hire_plan.xlsx"
        if fp.exists():
            try:
                xls = pd.ExcelFile(fp)
                if "hire_plan" in xls.sheet_names:
                    df_plan = xls.parse("hire_plan"); st.dataframe(df_plan, use_container_width=True, hide_index=True)
                    if "role" in df_plan and "hire_need" in df_plan: st.bar_chart(df_plan.set_index("role")["hire_need"])
                if "meta" in xls.sheet_names:
                    with st.expander(_("Hiring Plan Parameters")): st.table(xls.parse("meta"))
            except Exception as e: st.error(f"hire_plan.xlsx 表示エラー: {e}")
        else: st.info(_("Hire Plan") + " (hire_plan.xlsx) " + _("が見つかりません。"))

def display_leave_analysis_tab(tab_container, results_dict: dict | None = None):
    """Render leave analysis results using in-memory DataFrames."""
    results_dict = results_dict or {}
    with tab_container:
        st.subheader(_("Leave Analysis"))
        if not results_dict:
            st.info("No leave analysis results available.")
            return

        daily_df = results_dict.get("daily_leave_df")
        if isinstance(daily_df, pd.DataFrame) and not daily_df.empty:
            st.markdown("**日次・職員別取得データ**")
            st.dataframe(daily_df, use_container_width=True, hide_index=True)

        def _bar_chart(df: pd.DataFrame, title: str):
            if "period_unit" in df.columns and "total_leave_days" in df.columns:
                fig = px.bar(df, x="period_unit", y="total_leave_days", title=title)
                st.plotly_chart(fig, use_container_width=True)
            with st.expander("Data"):
                st.dataframe(df, use_container_width=True, hide_index=True)

        df_req_dow = results_dict.get("summary_dow_requested")
        if isinstance(df_req_dow, pd.DataFrame) and not df_req_dow.empty:
            _bar_chart(df_req_dow, "希望休 曜日別")

        df_req_mp = results_dict.get("summary_month_period_requested")
        if isinstance(df_req_mp, pd.DataFrame) and not df_req_mp.empty:
            _bar_chart(df_req_mp, "希望休 月内区分別")

        df_req_month = results_dict.get("summary_month_requested")
        if isinstance(df_req_month, pd.DataFrame) and not df_req_month.empty:
            _bar_chart(df_req_month, "希望休 月別")

        df_conc = results_dict.get("concentration_requested")
        if isinstance(df_conc, pd.DataFrame) and not df_conc.empty:
            st.markdown("**希望休 集中日判定**")
            st.dataframe(df_conc, use_container_width=True, hide_index=True)

        df_paid_dow = results_dict.get("summary_dow_paid")
        if isinstance(df_paid_dow, pd.DataFrame) and not df_paid_dow.empty:
            _bar_chart(df_paid_dow, "有給 曜日別")

        df_paid_month = results_dict.get("summary_month_paid")
        if isinstance(df_paid_month, pd.DataFrame) and not df_paid_month.empty:
            _bar_chart(df_paid_month, "有給 月別")

        df_balance = results_dict.get("staff_balance_daily")
        if isinstance(df_balance, pd.DataFrame) and not df_balance.empty:
            st.markdown("**勤務予定人数と希望休取得者数**")
            fig = px.line(df_balance, x="date", y=["total_staff", "leave_applicants_count", "non_leave_staff"],
                           markers=True)
            st.plotly_chart(fig, use_container_width=True)
            with st.expander("Data"):
                st.dataframe(df_balance, use_container_width=True, hide_index=True)

        df_staff = results_dict.get("staff_leave_list")
        if isinstance(df_staff, pd.DataFrame) and not df_staff.empty:
            st.markdown("**職員別休暇リスト**")
            st.dataframe(df_staff, use_container_width=True, hide_index=True)

def display_ppt_tab(tab_container, data_dir_ignored, key_prefix: str = ""):
    with tab_container:
        st.subheader(_("PPT Report"))
        button_key = f"dash_generate_ppt_button_{key_prefix or 'default'}"
        if st.button(_("Generate PowerPoint Report (β)"), key=button_key, use_container_width=True):
            st.info(_("Generating PowerPoint report..."))
            try:
                from pptx import Presentation 
                prs = Presentation()
                prs.slides.add_slide(prs.slide_layouts[5]).shapes.title.text = "ダッシュボードからのレポート"
                with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp_ppt_dash:
                    temp_ppt_dash_path = tmp_ppt_dash.name
                prs.save(temp_ppt_dash_path)
                with open(temp_ppt_dash_path, "rb") as ppt_file_data_dash:
                    st.download_button(
                        label=_("Download Report (PPTX)"), data=ppt_file_data_dash,
                        file_name=f"ShiftSuite_Dashboard_Report_{dt.datetime.now().strftime('%Y%m%d_%H%M')}.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        use_container_width=True
                    )
                Path(temp_ppt_dash_path).unlink()
                st.success(_("PowerPoint report ready."))
            except ImportError: st.error(_("python-pptx library required for PPT"))
            except Exception as e_ppt_dash: st.error(_("Error generating PowerPoint report") + f": {e_ppt_dash}")
        else:
            st.markdown(_("Click button to generate report."))
# Multi-file results display
if st.session_state.get("analysis_done", False) and st.session_state.analysis_results:
    st.divider()
    file_tabs = st.tabs(list(st.session_state.analysis_results.keys()))
    for tab_obj, fname in zip(file_tabs, st.session_state.analysis_results.keys()):
        with tab_obj:
            results = st.session_state.analysis_results[fname]
            st.subheader(f"Results for {fname}")
            data_dir = Path(results["out_dir_path_str"])
            tab_keys_en_dash = [
                "Overview", "Heatmap", "Shortage", "Fatigue", "Forecast",
                "Fairness", "Leave Analysis", "Cost Sim", "Hire Plan", "PPT Report"
            ]
            tab_labels_dash = [_(key) for key in tab_keys_en_dash]
            inner_tabs = st.tabs(tab_labels_dash)
            tab_func_map_dash = {
                "Overview": display_overview_tab,
                "Heatmap": display_heatmap_tab,
                "Shortage": display_shortage_tab,
                "Fatigue": display_fatigue_tab,
                "Forecast": display_forecast_tab,
                "Fairness": display_fairness_tab,
                "Leave Analysis": display_leave_analysis_tab,
                "Cost Sim": display_costsim_tab,
                "Hire Plan": display_hireplan_tab,
                "PPT Report": display_ppt_tab,
            }
            for i, key in enumerate(tab_keys_en_dash):
                if key in tab_func_map_dash:
                    if key == "PPT Report":
                        tab_func_map_dash[key](inner_tabs[i], data_dir, key_prefix=fname)
                    elif key == "Leave Analysis":
                        tab_func_map_dash[key](inner_tabs[i], results.get("leave_analysis_results", {}))
                    else:
                        tab_func_map_dash[key](inner_tabs[i], data_dir)



st.divider()
st.header(_("Dashboard (Upload ZIP)"))
zip_file_uploaded_dash_final_v3_display_main_dash = st.file_uploader(_("Upload ZIP file of 'out' folder"), type=["zip"], key="dashboard_zip_uploader_widget_final_v3_key_dash")

if zip_file_uploaded_dash_final_v3_display_main_dash:
    dashboard_temp_base = Path(tempfile.gettempdir()) / "ShiftSuite_Dash_Uploads"
    dashboard_temp_base.mkdir(parents=True, exist_ok=True)
    current_dash_tmp_dir = Path(tempfile.mkdtemp(prefix="dash_", dir=dashboard_temp_base))
    log.info(f"ダッシュボード用一時ディレクトリを作成: {current_dash_tmp_dir}")
    extracted_data_dir: Optional[Path] = None
    try:
        with zipfile.ZipFile(io.BytesIO(zip_file_uploaded_dash_final_v3_display_main_dash.read())) as zf:
            zf.extractall(current_dash_tmp_dir)
            if (current_dash_tmp_dir / "out").exists() and (current_dash_tmp_dir / "out" / "heat_ALL.xlsx").exists():
                extracted_data_dir = current_dash_tmp_dir / "out"
            elif (current_dash_tmp_dir / "heat_ALL.xlsx").exists():
                 extracted_data_dir = current_dash_tmp_dir
            else:
                found_heat_all = list(current_dash_tmp_dir.rglob("heat_ALL.xlsx"))
                if found_heat_all: 
                    extracted_data_dir = found_heat_all[0].parent
                else: 
                    st.error(_("heat_ALL.xlsx not found in ZIP"))
                    log.error(f"ZIP展開後、heat_ALL.xlsx が見つかりません in {current_dash_tmp_dir}")
                    st.stop()
        log.info(f"ダッシュボード表示用のデータディレクトリ: {extracted_data_dir}")
    except Exception as e_zip:
        st.error(_("Error during ZIP file extraction") + f": {e_zip}")
        log.error(f"ZIP展開中エラー: {e_zip}", exc_info=True)
        st.stop()

    import plotly.express as px
    import plotly.graph_objects as go

    tab_keys_en_dash = ["Overview", "Heatmap", "Shortage", "Fatigue", "Forecast", "Fairness", "Leave Analysis", "Cost Sim", "Hire Plan", "PPT Report"]
    tab_labels_dash = [_(key) for key in tab_keys_en_dash]
    tabs_obj_dash = st.tabs(tab_labels_dash)


    if extracted_data_dir:
        tab_function_map_dash = {
            "Overview": display_overview_tab,
            "Heatmap": display_heatmap_tab,
            "Shortage": display_shortage_tab,
            "Fatigue": display_fatigue_tab,
            "Forecast": display_forecast_tab,
            "Fairness": display_fairness_tab,
            "Leave Analysis": display_leave_analysis_tab,
            "Cost Sim": display_costsim_tab,
            "Hire Plan": display_hireplan_tab,
            "PPT Report": display_ppt_tab,
        }
        
        # 各タブに対応する表示関数を呼び出す
        for i, tab_key in enumerate(tab_keys_en_dash):
            if tab_key in tab_function_map_dash:
                if tab_key == "PPT Report":
                    tab_function_map_dash[tab_key](tabs_obj_dash[i], extracted_data_dir, key_prefix="zip")
                elif tab_key == "Leave Analysis":
                    tab_function_map_dash[tab_key](tabs_obj_dash[i], {})
                else:
                    tab_function_map_dash[tab_key](tabs_obj_dash[i], extracted_data_dir)
            else:
                with tabs_obj_dash[i]: 
                    st.subheader(_(tab_key))
                    st.info(f"{_(tab_key)} の表示は現在準備中です。")
    else:
        st.warning("ダッシュボードを表示するためのデータがロードされていません。ZIPファイルをアップロードしてください。")

if __name__ == "__main__" and not st_runtime_exists():
    import argparse
    log.info("CLIモードでapp.pyを実行します。")
    parser = argparse.ArgumentParser(description="Shift-Suite CLI (app.py経由のデバッグ用)")
    parser.add_argument("xlsx_file_cli", help="Excel シフト原本 (.xlsx)")
    parser.add_argument("--sheets_cli", nargs="+", required=True, help="解析対象のシート名")
    parser.add_argument("--header_cli", type=int, default=3, help="ヘッダー開始行 (1-indexed)")
    try:
        cli_args = parser.parse_args()
        log.info(f"CLI Args: file='{cli_args.xlsx_file_cli}', sheets={cli_args.sheets_cli}, header={cli_args.header_cli}")
    except SystemExit: 
        pass
    except Exception as e_cli:
        log.error(f"CLIモードでの実行中にエラー: {e_cli}", exc_info=True)
